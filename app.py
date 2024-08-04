import os
from langchain_google_genai import GoogleGenerativeAI
import streamlit as st
import base64
from pptx import Presentation
from pptx.util import Pt
from dotenv import load_dotenv

load_dotenv()

# Function work
GEMINI_API_KEY = os.getenv('API_KEY')
model = GoogleGenerativeAI(model='models/gemini-1.5-pro', google_api_key=GEMINI_API_KEY)
Title_Font_Size = Pt(30)
Slide_Font_Size = Pt(16)

def get_slide_title(topic, n):
    prompt = f"Generate {n} slide titles for the given topic '{topic}'"
    response = model.generate(prompts=[prompt], temperature=0.5)
    return response.generations[0][0].text.split("\n")

def get_slide_content(slide_title):
    prompt = f"Generate content on the given topic '{slide_title}'"
    response = model.generate(prompts=[prompt], temperature=0.5)
    return response.generations[0][0].text

def create_presentation(topic, slide_titles, slide_contents):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide_layout = prs.slide_layouts[1]

    # Title slide
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_slide.shapes.title.text = topic

    # Content slides
    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.placeholders[1].text = slide_content

        # Customize font size
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Slide_Font_Size

    prs.save(f"{topic}_presentation.pptx")

def get_ppt_download(topic):
    ppt_filename = f"{topic}_presentation.pptx"
    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the PowerPoint Presentation</a>'

# Streamlit app
st.title("PPT Generation from prompt:")
st.subheader("What do you want your presentation to be about?")
topic = st.text_input("Enter the topic here")
n = st.number_input("How many slides?", min_value=1, step=1)
generate_button = st.button("Generate PPT")

if generate_button and topic and n:
    st.info("Generating the presentation... Please wait.")
    slide_titles = get_slide_title(topic=topic, n=int(n))
    filtered_slide_titles = [title for title in slide_titles if title.strip()]
    slide_contents = [get_slide_content(title) for title in filtered_slide_titles]
    create_presentation(topic=topic, slide_titles=filtered_slide_titles, slide_contents=slide_contents)
    st.success("Presentation Generated Successfully")
    st.markdown(get_ppt_download(topic), unsafe_allow_html=True)
